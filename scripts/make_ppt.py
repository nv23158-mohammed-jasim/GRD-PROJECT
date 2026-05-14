from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Colors (LAB fitness app aesthetic: red/black/white)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK = RGBColor(0x18, 0x18, 0x1B)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
LIGHT_GRAY = RGBColor(0xE4, 0xE4, 0xE7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

def add_bg(slide, color=BLACK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_accent_bar(slide):
    # left red accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

def add_text(slide, text, left, top, width, height, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb

def add_bullets(slide, items, left, top, width, height, size=18, color=WHITE,
                bullet_color=RED, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r1.font.name = "Calibri"
        # text
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Calibri"
    return tb

def add_section_header(slide, label, title):
    add_accent_bar(slide)
    add_text(slide, label, Inches(0.6), Inches(0.45), Inches(8), Inches(0.4),
             size=14, bold=True, color=RED_LIGHT)
    add_text(slide, title, Inches(0.6), Inches(0.85), Inches(12), Inches(0.8),
             size=34, bold=True, color=WHITE)
    # underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.7),
                                   Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

def add_footer(slide, page):
    add_text(slide, "LAB — AI Fitness Tracker", Inches(0.6), Inches(7.05),
             Inches(6), Inches(0.3), size=10, color=GRAY)
    add_text(slide, f"Page {page}", Inches(12.2), Inches(7.05),
             Inches(0.8), Inches(0.3), size=10, color=GRAY, align=PP_ALIGN.RIGHT)

blank = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, BLACK)

# Decorative red gradient strip
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), SW, Inches(0.08))
strip.fill.solid(); strip.fill.fore_color.rgb = RED; strip.line.fill.background()

# Big LAB logo-ish title
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "LAB"
r.font.size = Pt(140); r.font.bold = True; r.font.color.rgb = RED; r.font.name = "Calibri"

add_text(s, "AI-Powered Fitness Tracking & Pose Detection",
         Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Graduation Project Report  •  2026",
         Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
         size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Department badge
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(4.4), Inches(4.9), Inches(4.5), Inches(0.55))
badge.fill.solid(); badge.fill.fore_color.rgb = DARK
badge.line.color.rgb = RED; badge.line.width = Pt(1.5)
tf = badge.text_frame; tf.margin_top = Inches(0.05)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "IT  •  CLOUD COMPUTING  •  12.CCP"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

add_text(s, "Presented by Class 12.CCP",
         Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Team
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "01  ·  THE TEAM", "Student Names & Roles")

team = [
    ("Mohammed Jasim", "Team Leader · Technical Writer · AI Model Designer"),
    ("Nasser Kuwaitan", "UI Designer"),
    ("Ali Ramzi", "Full Website Tester"),
    ("Mohammd Waleed", "Cloud Architect"),
    ("Ali Salman", "Cloud Architect"),
]

start_y = Inches(2.05)
card_h = Inches(0.85)
gap = Inches(0.1)
for i, (name, role) in enumerate(team):
    top = start_y + (card_h + gap) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.8), top, Inches(11.7), card_h)
    card.fill.solid(); card.fill.fore_color.rgb = DARK
    card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46); card.line.width = Pt(0.75)

    # number circle
    num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(1.1), top + Inches(0.16),
                             Inches(0.55), Inches(0.55))
    num.fill.solid(); num.fill.fore_color.rgb = RED; num.line.fill.background()
    tf = num.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1)
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE

    add_text(s, name, Inches(1.95), top + Inches(0.1),
             Inches(4.5), Inches(0.4), size=20, bold=True, color=WHITE)
    add_text(s, role, Inches(1.95), top + Inches(0.45),
             Inches(9.5), Inches(0.35), size=13, color=GRAY)

add_footer(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Department & Supervisors
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "02  ·  ACADEMIC INFO", "Department & Supervisors")

# Left: Department card
dep = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.5))
dep.fill.solid(); dep.fill.fore_color.rgb = DARK
dep.line.color.rgb = RED; dep.line.width = Pt(1.5)

add_text(s, "DEPARTMENT", Inches(1.05), Inches(2.35), Inches(5), Inches(0.35),
         size=12, bold=True, color=RED_LIGHT)
add_text(s, "Information Technology", Inches(1.05), Inches(2.7), Inches(5.3), Inches(0.5),
         size=22, bold=True, color=WHITE)

add_text(s, "SPECIALIZATION", Inches(1.05), Inches(3.6), Inches(5), Inches(0.35),
         size=12, bold=True, color=RED_LIGHT)
add_text(s, "Cloud Computing", Inches(1.05), Inches(3.95), Inches(5.3), Inches(0.5),
         size=22, bold=True, color=WHITE)

add_text(s, "CLASS", Inches(1.05), Inches(4.85), Inches(5), Inches(0.35),
         size=12, bold=True, color=RED_LIGHT)
add_text(s, "12.CCP", Inches(1.05), Inches(5.2), Inches(5.3), Inches(0.5),
         size=22, bold=True, color=WHITE)

# Right: Supervisors card
sup = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(6.8), Inches(2.1), Inches(5.7), Inches(4.5))
sup.fill.solid(); sup.fill.fore_color.rgb = DARK
sup.line.color.rgb = RED; sup.line.width = Pt(1.5)

add_text(s, "SUPERVISING TEACHERS", Inches(7.05), Inches(2.35),
         Inches(5.3), Inches(0.4), size=12, bold=True, color=RED_LIGHT)

supervisors = [
    "Mr. Khalid Maroof",
    "Mr. Abdullrahman Khaled",
    "Dr. Raghu Dumpati",
]
for i, name in enumerate(supervisors):
    top = Inches(2.95) + Inches(1.15) * i
    # bullet square
    sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(7.05), top + Inches(0.18), Inches(0.12), Inches(0.45))
    sq.fill.solid(); sq.fill.fore_color.rgb = RED; sq.line.fill.background()
    add_text(s, name, Inches(7.4), top + Inches(0.1),
             Inches(5), Inches(0.55), size=20, bold=True, color=WHITE)
    add_text(s, "Project Supervisor", Inches(7.4), top + Inches(0.55),
             Inches(5), Inches(0.35), size=12, color=GRAY)

add_footer(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Brief Description
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "03  ·  ABOUT THE PROJECT", "Brief Description")

desc = ("LAB is a full-stack web application that turns any device with a camera "
        "into a personal AI fitness coach. It uses TensorFlow.js and the MoveNet "
        "pose-detection model to track the user's body in real time, count "
        "repetitions, and grade exercise form — no wearables or sensors required.")
add_text(s, desc, Inches(0.8), Inches(2.05), Inches(11.7), Inches(1.5),
         size=17, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# Three feature cards
features = [
    ("Workout Tracking",
     "Push-ups, squats, plank with live rep counting and form grading."),
    ("Neon Run Game",
     "Side-scrolling game controlled by jogging, jumping, and push-ups."),
    ("Boxing Mode",
     "Shadow-boxing trainer with punches, dodges, and blocks detected by camera."),
]
card_w = Inches(3.9); card_h = Inches(2.5); top = Inches(4.0)
for i, (title, body) in enumerate(features):
    left = Inches(0.8) + (card_w + Inches(0.1)) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.fill.solid(); card.fill.fore_color.rgb = DARK
    card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46); card.line.width = Pt(0.75)
    # top accent
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.08))
    acc.fill.solid(); acc.fill.fore_color.rgb = RED; acc.line.fill.background()

    add_text(s, title, left + Inches(0.25), top + Inches(0.3),
             card_w - Inches(0.5), Inches(0.5), size=18, bold=True, color=WHITE)
    add_text(s, body, left + Inches(0.25), top + Inches(0.95),
             card_w - Inches(0.5), Inches(1.4), size=13, color=GRAY)

add_footer(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Purpose & Beneficiaries
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "04  ·  WHY WE BUILT IT", "Purpose & Target Beneficiaries")

# Purpose
add_text(s, "PURPOSE", Inches(0.8), Inches(2.05), Inches(6), Inches(0.4),
         size=13, bold=True, color=RED_LIGHT)
purpose_items = [
    "Make professional fitness coaching accessible to everyone, for free.",
    "Replace expensive wearables and personal trainers with a phone or laptop camera.",
    "Use AI to give real-time feedback on form, posture, and rep count.",
    "Make exercise fun through gamification (Neon Run, Boxing Mode).",
]
add_bullets(s, purpose_items, Inches(0.8), Inches(2.5),
            Inches(11.7), Inches(2.5), size=15)

# Beneficiaries
add_text(s, "TARGET BENEFICIARIES", Inches(0.8), Inches(4.95),
         Inches(6), Inches(0.4), size=13, bold=True, color=RED_LIGHT)

beneficiaries = [
    ("Home Users", "People who can't afford gyms"),
    ("Students", "Quick fitness breaks"),
    ("Beginners", "Need form guidance"),
    ("Gamers", "Active gameplay"),
]
card_w = Inches(2.85); card_h = Inches(1.55); top = Inches(5.4)
for i, (title, sub) in enumerate(beneficiaries):
    left = Inches(0.8) + (card_w + Inches(0.07)) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.fill.solid(); card.fill.fore_color.rgb = DARK
    card.line.color.rgb = RED; card.line.width = Pt(1)
    add_text(s, title, left + Inches(0.2), top + Inches(0.25),
             card_w - Inches(0.4), Inches(0.5), size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, sub, left + Inches(0.2), top + Inches(0.85),
             card_w - Inches(0.4), Inches(0.5), size=12, color=GRAY,
             align=PP_ALIGN.CENTER)

add_footer(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Current Stage
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "05  ·  STATUS", "Current Development Stage")

# Big status badge
status = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.8), Inches(2.05), Inches(11.7), Inches(1.0))
status.fill.solid(); status.fill.fore_color.rgb = RED
status.line.fill.background()
tf = status.text_frame; tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "✓  FULLY DEPLOYED & LIVE IN PRODUCTION"
r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE

# Tech stack grid
add_text(s, "WHAT'S COMPLETE", Inches(0.8), Inches(3.35),
         Inches(6), Inches(0.4), size=13, bold=True, color=RED_LIGHT)

complete = [
    "Full React + TypeScript frontend with dark fitness theme",
    "Node.js + Express REST API backend with type-safe contracts",
    "PostgreSQL database (Neon) with Drizzle ORM",
    "Google OAuth, Microsoft OAuth, and Email/Password login (JWT-based)",
    "Real-time pose detection for push-ups, squats, and plank",
    "Neon Run game with 5 unlockable stages and 3 enemy types",
    "Boxing Mode with rounds, voice commands, and scoring",
    "Admin panel with user management, audit logs, and search",
    "Production deployment: Render (backend) + AWS Amplify (frontend)",
]
add_bullets(s, complete, Inches(0.8), Inches(3.8),
            Inches(11.7), Inches(3.2), size=14, line_spacing=1.15)

add_footer(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Challenges
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "06  ·  WHAT WE LEARNED", "Challenges Faced During Development")

challenges = [
    ("Cross-site cookies blocked on iOS Safari",
     "Frontend on Amplify and backend on Render are different domains. Safari blocked our session cookies. We solved it by switching authentication to JWT tokens stored in localStorage and sent as Authorization headers."),
    ("Pose detection accuracy & calibration",
     "Tuning the MoveNet model thresholds for push-up depth, squat angles, and plank alignment took many iterations to feel natural across different body types and camera angles."),
    ("Cloud cold starts",
     "Render's free tier puts the backend to sleep after 15 minutes of inactivity. First request after sleep takes 30+ seconds — we documented this and explored keep-alive strategies."),
    ("Multi-platform deployment coordination",
     "Coordinating GitHub → Render (backend) and GitHub → Amplify (frontend) auto-deploys, plus database migrations on Neon, required careful environment-variable management across three platforms."),
    ("Real-time game performance",
     "Running TensorFlow.js pose detection at 30+ FPS while rendering a side-scrolling game in the same browser tab pushed performance limits. Required careful optimization of the render loop."),
]

top = Inches(2.05)
for i, (title, body) in enumerate(challenges):
    h = Inches(0.95)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.8), top, Inches(11.7), h)
    card.fill.solid(); card.fill.fore_color.rgb = DARK
    card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46); card.line.width = Pt(0.5)
    # left accent
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.08), h)
    acc.fill.solid(); acc.fill.fore_color.rgb = RED; acc.line.fill.background()

    add_text(s, title, Inches(1.05), top + Inches(0.1),
             Inches(11.3), Inches(0.4), size=14, bold=True, color=WHITE)
    add_text(s, body, Inches(1.05), top + Inches(0.45),
             Inches(11.3), Inches(0.5), size=11, color=GRAY)
    top += h + Inches(0.05)

add_footer(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Future Plans
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "07  ·  WHAT'S NEXT", "Future Development Plans")

plans = [
    ("More Exercises",
     "Add lunges, burpees, mountain climbers, and yoga poses to the pose detection library."),
    ("Native Mobile App",
     "Build a React Native / Expo version for iOS and Android with offline workout caching."),
    ("Social & Leaderboards",
     "Friend challenges, weekly leaderboards, and sharable workout summaries."),
    ("AI Personal Trainer",
     "GPT-powered coach that adapts difficulty, builds custom plans, and gives voice cues."),
    ("Monetization",
     "Premium tier via RevenueCat with advanced analytics, custom programs, and meal plans."),
    ("Wearable Integration",
     "Sync with Apple Watch and Google Fit for heart-rate and calorie accuracy."),
]

card_w = Inches(5.8); card_h = Inches(1.55)
positions = [(0, 0), (1, 0), (0, 1), (1, 1), (0, 2), (1, 2)]
start_top = Inches(2.05); start_left = Inches(0.8)
gap_x = Inches(0.1); gap_y = Inches(0.1)
for i, (title, body) in enumerate(plans):
    col, row = positions[i]
    left = start_left + (card_w + gap_x) * col
    top = start_top + (card_h + gap_y) * row
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.fill.solid(); card.fill.fore_color.rgb = DARK
    card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46); card.line.width = Pt(0.75)
    # number badge
    num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             left + Inches(0.25), top + Inches(0.25),
                             Inches(0.45), Inches(0.45))
    num.fill.solid(); num.fill.fore_color.rgb = RED; num.line.fill.background()
    tf = num.text_frame; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1)
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

    add_text(s, title, left + Inches(0.85), top + Inches(0.25),
             card_w - Inches(1.1), Inches(0.45), size=16, bold=True, color=WHITE)
    add_text(s, body, left + Inches(0.85), top + Inches(0.7),
             card_w - Inches(1.1), Inches(0.8), size=11, color=GRAY)

add_footer(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Prototype / Demo
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "08  ·  TRY IT NOW", "Prototype & Live Demonstration")

# Big availability badge
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.85))
badge.fill.solid(); badge.fill.fore_color.rgb = RED; badge.line.fill.background()
tf = badge.text_frame; tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "PROTOTYPE AVAILABLE — LIVE & PUBLICLY ACCESSIBLE"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE

# URL cards
urls = [
    ("Live Web App (Frontend)",
     "main.dfi1vofcktrh3.amplifyapp.com",
     "Hosted on AWS Amplify"),
    ("Backend API",
     "grd-project-server.onrender.com",
     "Hosted on Render"),
    ("Source Code",
     "github.com/nv23158-mohammed-jasim/GRD-PROJECT",
     "Full source on GitHub"),
]

top = Inches(3.3)
for i, (title, url, sub) in enumerate(urls):
    h = Inches(1.15)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.8), top, Inches(11.7), h)
    card.fill.solid(); card.fill.fore_color.rgb = DARK
    card.line.color.rgb = RED; card.line.width = Pt(1)

    add_text(s, title, Inches(1.1), top + Inches(0.15),
             Inches(5), Inches(0.4), size=14, bold=True, color=RED_LIGHT)
    add_text(s, url, Inches(1.1), top + Inches(0.5),
             Inches(11), Inches(0.4), size=18, bold=True, color=WHITE)
    add_text(s, sub, Inches(1.1), top + Inches(0.85),
             Inches(11), Inches(0.3), size=11, color=GRAY)
    top += h + Inches(0.1)

add_footer(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Project Photos placeholder
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)
add_section_header(s, "09  ·  GALLERY", "Project Photos")

# Try to embed up to 4 screenshot placeholders
photo_paths = [
    "attached_assets/image_1776638568991.png",
    "attached_assets/image_1776634602123.png",
    "attached_assets/image_1776633549672.png",
    "attached_assets/image_1776633503273.png",
]
photo_paths = [p for p in photo_paths if os.path.exists(p)]

if photo_paths:
    cols = 2
    cell_w = Inches(5.9); cell_h = Inches(2.35)
    start_left = Inches(0.8); start_top = Inches(2.1)
    gap_x = Inches(0.1); gap_y = Inches(0.15)
    for i, path in enumerate(photo_paths[:4]):
        col = i % cols; row = i // cols
        left = start_left + (cell_w + gap_x) * col
        top = start_top + (cell_h + gap_y) * row
        # frame
        frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, cell_w, cell_h)
        frame.fill.solid(); frame.fill.fore_color.rgb = DARK
        frame.line.color.rgb = RED; frame.line.width = Pt(1)
        # image fitted inside
        try:
            s.shapes.add_picture(path, left + Inches(0.05), top + Inches(0.05),
                                  width=cell_w - Inches(0.1), height=cell_h - Inches(0.1))
        except Exception:
            pass

    add_text(s, "Screenshots from the LAB project (database, admin panel, app UI)",
             Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)
else:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(2.5), Inches(2.5), Inches(8.3), Inches(3.5))
    box.fill.solid(); box.fill.fore_color.rgb = DARK
    box.line.color.rgb = RED; box.line.width = Pt(1.5)
    add_text(s, "📷", Inches(2.5), Inches(3.0),
             Inches(8.3), Inches(0.8), size=48, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, "Add Screenshots Here",
             Inches(2.5), Inches(3.9), Inches(8.3), Inches(0.5),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Insert photos of the workout tracker, Neon Run game, Boxing Mode, and Admin Panel.",
             Inches(2.5), Inches(4.5), Inches(8.3), Inches(0.6),
             size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_footer(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Thank You
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); add_bg(s, BLACK)

# Big LAB
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(2))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "THANK YOU"
r.font.size = Pt(80); r.font.bold = True; r.font.color.rgb = RED

add_text(s, "Questions & Discussion",
         Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.6),
         size=26, color=WHITE, align=PP_ALIGN.CENTER)

# strip
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.6),
                            Inches(2.3), Inches(0.06))
strip.fill.solid(); strip.fill.fore_color.rgb = RED; strip.line.fill.background()

add_text(s, "LAB — AI-Powered Fitness Tracking",
         Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.4),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Class 12.CCP  •  Cloud Computing  •  2026",
         Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.4),
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

# Save
out = "attached_assets/LAB_Project_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
